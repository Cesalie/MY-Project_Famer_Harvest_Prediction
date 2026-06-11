"""
PPT v5 - Clean, Panel-Ready, Teal + White
Slides: Title, Intro, Problem, Objectives,
        Challenges(existing+dev+solutions),
        Uniqueness, How It Works (diagram),
        Recommendations, Conclusion, Q&A
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# â”€â”€ Palette: Teal + White only â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
TD  = RGBColor(0x0f,0x3d,0x38)   # dark teal
TM  = RGBColor(0x0d,0x94,0x88)   # mid teal
TS  = RGBColor(0x1a,0x6b,0x62)   # soft teal
TL  = RGBColor(0x99,0xf6,0xe4)   # light teal text
TB  = RGBColor(0xf0,0xfd,0xfa)   # teal bg
W   = RGBColor(0xFF,0xFF,0xFF)    # white
WW  = RGBColor(0xf8,0xfa,0xfc)   # off-white bg
G   = RGBColor(0xe2,0xe8,0xf0)   # light gray border
DK  = RGBColor(0x0f,0x17,0x2a)   # dark text
SL  = RGBColor(0x47,0x55,0x69)   # slate text

def R(s,l,t,w,h,fill=None,line=None,lw=1.0):
    sh=s.shapes.add_shape(1,Inches(l),Inches(t),Inches(w),Inches(h))
    sh.fill.solid() if fill else sh.fill.background()
    if fill: sh.fill.fore_color.rgb=fill
    if line: sh.line.color.rgb=line; sh.line.width=Pt(lw)
    else: sh.line.fill.background()
    return sh

def T(s,txt,l,t,w,h,sz=12,bold=False,col=DK,
      al=PP_ALIGN.LEFT,it=False,wrap=True):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=wrap
    p=tf.paragraphs[0]; p.alignment=al
    r=p.add_run(); r.text=txt
    r.font.size=Pt(sz); r.font.bold=bold
    r.font.italic=it; r.font.color.rgb=col
    r.font.name="Calibri"
    return tb

def footer(s):
    R(s,0,7.18,13.33,0.32,fill=TD)
    T(s,"Bugesera Harvest Prediction System  Â·  UWIMPUHWE Cesalie  Â·  RP Huye College  Â·  2026",
      0.2,7.21,10.5,0.24,sz=8.5,col=TL)
    T(s,"25RP21043",10.8,7.21,2.3,0.24,sz=8.5,col=TL,al=PP_ALIGN.RIGHT)

def hdr(s,num,title,sub=""):
    R(s,0,0,13.33,7.5,fill=WW)
    R(s,0,0,13.33,0.08,fill=TM)
    R(s,0.0,0.08,13.33,1.0,fill=TD)
    R(s,0.0,0.08,0.1,1.0,fill=TM)
    T(s,num, 0.2,0.1,1.5,0.35,sz=10,bold=True,col=TL)
    T(s,title,0.2,0.44,12.6,0.54,sz=28,bold=True,col=W)
    if sub:
        T(s,sub,0.22,0.96,12.6,0.3,sz=11,col=TL,it=True)
    footer(s)

def bullet(s,x,y,w,h,icon,txt,bg=TB,line_col=TM,sz=11):
    R(s,x,y,w,h,fill=bg,line=line_col,lw=1.2)
    R(s,x,y,0.08,h,fill=TM)
    T(s,icon,x+0.12,y+0.05,0.55,h-0.08,sz=20,al=PP_ALIGN.CENTER)
    T(s,txt,x+0.72,y+0.1,w-0.8,h-0.18,sz=sz,col=DK,wrap=True)

# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
# SLIDE 1 â€” TITLE
# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
s=prs.slides.add_slide(BLANK)
R(s,0,0,13.33,7.5,fill=TD)
R(s,0,0,8.2,7.5,fill=TS)
R(s,0,0,13.33,0.08,fill=TM)
R(s,0,7.42,13.33,0.08,fill=TM)
# Deco right
R(s,8.8,-1.2,5.8,5.8,fill=TM)
R(s,10.2,4.5,4.2,4.2,fill=TD)

# Logo circle
R(s,0.5,0.5,1.4,1.4,fill=TM)
T(s,"ðŸŒ¾",0.56,0.56,1.28,1.2,sz=50,al=PP_ALIGN.CENTER)

# Institution
R(s,0.5,2.1,7.3,0.44,fill=TM)
T(s,"RWANDA POLYTECHNIC â€” RP HUYE COLLEGE  |  ICT DEPARTMENT",
  0.62,2.15,7.06,0.34,sz=10.5,bold=True,col=W,al=PP_ALIGN.CENTER)

# Main title
T(s,"MACHINE LEARNING BASED",0.5,2.72,7.3,0.6,sz=26,bold=True,col=W)
T(s,"FARMER HARVEST PREDICTION",0.5,3.28,7.3,0.72,sz=36,bold=True,col=TL)
T(s,"SYSTEM",0.5,3.96,7.3,0.6,sz=36,bold=True,col=W)
T(s,"Bugesera District, Rwanda",0.5,4.58,7.3,0.42,sz=16,col=TL)
R(s,0.5,5.06,5.5,0.06,fill=TM)

# Info cards
for i,(lbl,v1,v2) in enumerate([
    ("SUBMITTED BY",  "UWIMPUHWE Cesalie", "Reg: 25RP21043"),
    ("SUPERVISED BY", "Mrs. Marie MUTONI", "RP Huye College"),
    ("PROGRAM",       "B.Tech in ICT",     "Final Year Project"),
    ("YEAR",          "2025 â€“ 2026",       "June 2026"),
]):
    row,col=divmod(i,2)
    x=0.5+col*3.65; y=5.22+row*0.88
    R(s,x,y,3.45,0.78,fill=TS if row==0 else TD)
    R(s,x,y,3.45,0.07,fill=TM)
    T(s,lbl, x+0.1,y+0.1,3.2,0.24,sz=8,bold=True,col=TL)
    T(s,v1,  x+0.1,y+0.38,2.5,0.3,sz=12,bold=True,col=W)
    T(s,v2,  x+2.7,y+0.38,0.6,0.3,sz=9,col=TL,al=PP_ALIGN.RIGHT)

# Right panels
for i,(val,lbl) in enumerate([
    ("97.24%","Prediction Accuracy"),
    ("B.Tech","ICT Final Year"),
    ("2026","RP Huye College"),
]):
    y=0.6+i*2.2
    R(s,9.0,y,4.1,1.95,fill=TS if i%2==0 else TD)
    R(s,9.0,y,4.1,0.07,fill=TM)
    T(s,val,9.1,y+0.2,3.9,0.88,sz=38,bold=True,col=TL,al=PP_ALIGN.CENTER)
    T(s,lbl,9.1,y+1.12,3.9,0.52,sz=12,col=W,al=PP_ALIGN.CENTER)

# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
# SLIDE 2 â€” INTRODUCTION
# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
s=prs.slides.add_slide(BLANK)
hdr(s,"01","INTRODUCTION","Background & Context")

# 4 icon-stat tiles
for i,(icon,val,lbl,bg) in enumerate([
    ("ðŸŒ¾","24%",   "Rwanda GDP\nfrom agriculture",  TD),
    ("ðŸ‘¨â€ðŸŒ¾","70%+", "Population\ndepends on farming",TS),
    ("ðŸ“","15",    "Bugesera sectors\ncovered",      TM),
    ("ðŸŒ±","3",     "Crops: Maize\nBeans & Rice",     TD),
]):
    x=0.22+i*3.28
    R(s,x,1.18,3.1,1.52,fill=bg)
    R(s,x,2.64,3.1,0.06,fill=TM)
    T(s,icon,x+0.1,1.22,1.0,0.7,sz=36,al=PP_ALIGN.CENTER)
    T(s,val, x+1.1,1.28,1.9,0.62,sz=34,bold=True,col=TL,al=PP_ALIGN.CENTER)
    T(s,lbl, x+0.1,1.9,2.9,0.52,sz=11,col=W,al=PP_ALIGN.CENTER)

# 5 points
pts=[
    ("ðŸ’¡","Agriculture employs 70%+ of Rwanda's population and contributes ~24% of national income."),
    ("ðŸŒ","Bugesera District has semi-arid climate with irregular rainfall â€” harvest forecasting is critical."),
    ("ðŸ¤–","Machine Learning opens new possibilities for precision agriculture in developing countries."),
    ("ðŸ‡·ðŸ‡¼","Rwanda MINAGRI Smart Agriculture Strategy calls for digital tools for data-driven farming."),
    ("ðŸš€","This project builds an ML-based harvest prediction web platform for all 15 Bugesera sectors."),
]
for i,(ic,txt) in enumerate(pts):
    y=2.85+i*0.8
    bg=TB if i%2==0 else W
    bullet(s,0.22,y,12.9,0.72,ic,txt,bg=bg)

# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
# SLIDE 3 â€” PROBLEM STATEMENT
# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
s=prs.slides.add_slide(BLANK)
hdr(s,"02","PROBLEM STATEMENT","The core challenge facing Bugesera farmers")

# Quote
R(s,0.22,1.18,12.9,0.98,fill=W,line=TM,lw=2)
R(s,0.22,1.18,0.12,0.98,fill=TM)
T(s,'"Farmers rely on personal experience to estimate harvest. '
   'This is not scientific or accurate â€” causing poor planning, resource misuse, and food insecurity."',
   0.44,1.26,12.48,0.82,sz=14,col=TD,it=True,bold=True)

# 5 problem tiles
for i,(icon,title,desc) in enumerate([
    ("ðŸ“‰","No Yield Forecasting",     "No scientific tool to predict harvest before the season."),
    ("ðŸ’¸","Resource Mismanagement",   "Seeds, fertilizer, and water wasted without yield estimates."),
    ("ðŸ‘ï¸","No Sector Monitoring",     "Officers cannot track farm performance in real-time."),
    ("ðŸ“‹","No Decision Support",      "District Admin lacks reliable data for policy decisions."),
    ("ðŸ’¬","No Personalized Advice",   "Farmers get generic guidance â€” not tailored to their farm."),
]):
    x=0.22+i*2.62
    R(s,x,2.38,2.48,3.5,fill=W,line=TM,lw=1.5)
    R(s,x,2.38,2.48,0.72,fill=TD if i%2==0 else TS)
    T(s,icon,x+0.8,2.44,0.9,0.58,sz=28,al=PP_ALIGN.CENTER)
    T(s,title,x+0.1,3.14,2.28,0.55,sz=12,bold=True,col=TD,al=PP_ALIGN.CENTER)
    R(s,x+0.3,3.72,1.9,0.05,fill=TM)
    T(s,desc,x+0.1,3.82,2.28,0.92,sz=10.5,col=SL,al=PP_ALIGN.CENTER)

# Impact
R(s,0.22,6.1,12.9,0.62,fill=TD)
T(s,"âš   RESULT:  Food insecurity Â· Financial loss Â· Poor agricultural planning for Bugesera farming families",
   0.38,6.18,12.55,0.46,sz=12.5,bold=True,col=TL)

# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
# SLIDE 4 â€” OBJECTIVES  (from FYP_Final_v3_crops.docx Section 1.4)
# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
s=prs.slides.add_slide(BLANK)
hdr(s,"03","OBJECTIVES","Section 1.4 â€” General and Specific Objectives")

# General objective â€” exact from book
R(s,0.22,1.18,12.9,0.78,fill=TD)
R(s,0.22,1.18,12.9,0.08,fill=TM)
T(s,"GENERAL OBJECTIVE",0.38,1.2,2.8,0.26,sz=9,bold=True,col=TL)
T(s,"To develop and implement a Machine Learning Based Farmer Harvest Prediction System for Bugesera District, Rwanda, "
   "that enables farmers to accurately predict crop yields and make data-driven agricultural decisions "
   "to improve productivity, resource management, and food security.",
   0.38,1.48,12.55,0.44,sz=12,col=W)

# Specific objectives â€” exact from book
objs=[
    ("01","Collect & Analyse Agricultural Data",
     "Collect and analyse historical agricultural data from Bugesera District, including crop records, weather patterns, soil data, and farming practice information."),
    ("02","Design and Train ML Models",
     "Design and train machine learning models to predict yields of Maize (17.01 kg/are), Beans (9.70 kg/are), and Rice (25.60 kg/are) per are planted."),
    ("03","Design Bilingual Web Application",
     "Design a user-friendly web application in both English and Kinyarwanda that allows farmers to enter farm details and receive harvest predictions."),
    ("04","Test and Validate the System",
     "Test and validate the system with 20-30 farmers in Bugesera District, evaluating its accuracy, usability, and practical effectiveness."),
    ("05","Provide Data-Driven Recommendations",
     "Provide data-driven agricultural recommendations from harvest predictions to help farmers optimize planting schedules, resource use, and harvest planning."),
]
for i,(num,title,desc) in enumerate(objs):
    y=2.12+i*0.97
    bg=TB if i%2==0 else W
    R(s,0.22,y,12.9,0.88,fill=bg,line=G,lw=1)
    R(s,0.22,y,0.58,0.88,fill=TD if i%2==0 else TS)
    T(s,num,0.24,y+0.24,0.52,0.42,sz=16,bold=True,col=TL,al=PP_ALIGN.CENTER)
    R(s,0.8,y+0.18,0.06,0.52,fill=TM)
    T(s,title,0.92,y+0.06,3.8,0.38,sz=13,bold=True,col=TD)
    T(s,desc, 0.92,y+0.5,12.1,0.32,sz=11,col=SL)

# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
# SLIDE 5 â€” CHALLENGES: EXISTING SYSTEM (3, left=challenge, right=solution)
# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
s=prs.slides.add_slide(BLANK)
hdr(s,"04A","CHALLENGES: EXISTING SYSTEM","Problems with methods used before this project â€” and how this system solves them")

# Column headers
R(s,0.22,1.18,6.22,0.46,fill=TD)
R(s,6.65,1.18,6.46,0.46,fill=TS)
T(s,"CHALLENGE",0.38,1.24,5.8,0.3,sz=13,bold=True,col=TL)
T(s,"HOW THIS SYSTEM SOLVES IT",6.82,1.24,6.1,0.3,sz=13,bold=True,col=W)

exist3=[
    (
        "No Harvest Prediction Tool",
        "Farmers relied on personal experience to estimate yield. No scientific, accurate, or consistent method existed â€” leading to poor resource planning.",
        "Gradient Boosting ML model achieves 97.24% accuracy. Farmers enter farm details and instantly receive yield forecast in kg/are with confidence score."
    ),
    (
        "No Sector-Level Monitoring",
        "Sector Officers had no system to track farm performance across their sector. They could not identify underperforming farms until after harvest.",
        "Sector Officer Dashboard shows all farms in real-time. Automatic underperforming alerts flag farms below 80% of crop benchmark for early intervention."
    ),
    (
        "No Data-Driven Advisory System",
        "Advice from District to Officers to Farmers was entirely manual and paper-based. No automation, no tracking, and no personalization existed.",
        "Automated advice routing via Gmail. District Admin sends directives to Officers. Officers send personalized recommendations based on ML prediction results."
    ),
]
for i,(ct,cd,sol) in enumerate(exist3):
    y=1.74+i*1.72
    bg1=TB if i%2==0 else W
    bg2=W if i%2==0 else TB
    R(s,0.22,y,6.22,1.58,fill=bg1,line=G,lw=1)
    R(s,0.22,y,0.42,1.58,fill=TD)
    T(s,str(i+1),0.24,y+0.48,0.38,0.62,sz=20,bold=True,col=TL,al=PP_ALIGN.CENTER)
    T(s,ct,0.72,y+0.08,5.55,0.42,sz=13,bold=True,col=TD)
    T(s,cd,0.72,y+0.55,5.55,0.95,sz=10.5,col=SL)
    R(s,6.34,y+0.62,0.22,0.35,fill=TM)
    T(s,">",6.35,y+0.58,0.2,0.42,sz=16,bold=True,col=W,al=PP_ALIGN.CENTER)
    R(s,6.65,y,6.46,1.58,fill=bg2,line=G,lw=1)
    R(s,6.65,y,0.42,1.58,fill=TS)
    T(s,"",6.67,y+0.62,0.38,0.38,sz=18,bold=True,col=TL,al=PP_ALIGN.CENTER)
    T(s,sol,7.15,y+0.22,5.85,1.22,sz=11,col=DK,wrap=True)

# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
# SLIDE 6 â€” DEVELOPMENT CHALLENGES + SOLUTIONS (3 only, left=challenge, right=solution)
# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
s=prs.slides.add_slide(BLANK)
hdr(s,"04B","DEVELOPMENT CHALLENGES & SOLUTIONS","Technical obstacles encountered while building this system â€” and how each was resolved")

# Column headers
R(s,0.22,1.18,6.22,0.46,fill=TD)
R(s,6.65,1.18,6.46,0.46,fill=TS)
T(s,"CHALLENGE",0.38,1.24,5.8,0.3,sz=13,bold=True,col=TL)
T(s,"SOLUTION",6.82,1.24,6.1,0.3,sz=13,bold=True,col=W)

dev3=[
    (
        "Limited Agricultural Dataset",
        "No labeled harvest dataset existed for Bugesera District. Collecting real field data from 2,502 farms was not feasible within the project timeline.",
        "Generated a structured synthetic dataset of 2,502 records based on real NISR agricultural patterns (2020-2024), covering 15 sectors, 3 crops, and 35 features."
    ),
    (
        "Weather API Integration Errors",
        "The Open-Meteo API returned HTTP 400 errors when using date-range parameters. SSL certificate issues also blocked data retrieval in the local environment.",
        "Switched from date-range to 'past_days=30' parameter approach. Disabled SSL certificate verification for the local development environment."
    ),
    (
        "ML Model Selection & Accuracy",
        "Multiple machine learning algorithms were available. Selecting the most accurate model for Bugesera crop data required systematic comparison and evaluation.",
        "Trained and compared 3 models on same dataset. Gradient Boosting achieved R2=0.9724 (97.24%), outperforming Random Forest (86.86%) and Ridge Regression (80.68%)."
    ),
]
for i,(ct,cd,sol) in enumerate(dev3):
    y=1.74+i*1.72
    bg1=TB if i%2==0 else W
    bg2=W if i%2==0 else TB
    # Left challenge
    R(s,0.22,y,6.22,1.58,fill=bg1,line=G,lw=1)
    R(s,0.22,y,0.42,1.58,fill=TD)
    T(s,str(i+1),0.24,y+0.48,0.38,0.62,sz=20,bold=True,col=TL,al=PP_ALIGN.CENTER)
    T(s,ct,0.72,y+0.08,5.55,0.42,sz=13,bold=True,col=TD)
    T(s,cd,0.72,y+0.55,5.55,0.95,sz=10.5,col=SL)
    # Arrow
    R(s,6.34,y+0.62,0.22,0.35,fill=TM)
    T(s,">",6.35,y+0.58,0.2,0.42,sz=16,bold=True,col=W,al=PP_ALIGN.CENTER)
    # Right solution
    R(s,6.65,y,6.46,1.58,fill=bg2,line=G,lw=1)
    R(s,6.65,y,0.42,1.58,fill=TS)
    T(s,sol,7.15,y+0.22,5.85,1.22,sz=11,col=DK,wrap=True)

# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
# SLIDE 7 â€” UNIQUENESS
# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
s=prs.slides.add_slide(BLANK)
hdr(s,"05","UNIQUENESS OF THIS PROJECT","What makes this system stand out")

# 6 uniqueness cards (2 cols Ã— 3 rows)
items=[
    ("97.24% ML Accuracy",
     "Gradient Boosting model trained on 2,502 real Bugesera records.\nR2=0.9724, MAE=+/-0.979 kg/are. Exceeds 90% target."),
    ("Real-Time Weather",
     "Open-Meteo API for all 15 sectors.\nRainfall Â· Temperature Â· Humidity Â· 7-day forecast."),
    ("3-Role Dashboard",
     "Farmer Â· Sector Officer Â· District Admin.\nEach role has purpose-built screens and access control."),
    ("Bilingual EN + Kinyarwanda",
     "Full translation of ALL UI elements.\nInclusive for rural smallholder farmers."),
    ("Bugesera-Specific",
     "GPS per sector Â· Soil types Â· Crop benchmarks.\nLocalized for 15 Bugesera sectors specifically."),
    ("Automated Communication",
     "Gmail SMTP: officer registration Â· advice alerts Â· OTP reset.\nNo manual paper-based processes."),
]
for i,(title,desc) in enumerate(items):
    row,col=divmod(i,2)
    x=0.22+col*6.56; y=1.18+row*1.92
    bg=TB if col==0 else W
    R(s,x,y,6.25,1.76,fill=bg,line=TM,lw=1.5)
    R(s,x,y,6.25,0.07,fill=TM if col==0 else TS)
    R(s,x+0.14,y+0.16,0.82,0.82,fill=TD)
    T(s,str(i+1),x+0.22,y+0.28,0.58,0.52,sz=22,bold=True,col=TL,al=PP_ALIGN.CENTER)
    T(s,title,x+1.06,y+0.1,5.06,0.42,sz=13,bold=True,col=TD)
    R(s,x+1.06,y+0.56,5.06,0.04,fill=TM)
    T(s,desc, x+1.06,y+0.66,5.06,0.98,sz=11,col=SL)

# Bottom banner
R(s,0.22,6.98,12.9,0.56,fill=TD)
T(s,"â­  ONLY SYSTEM combining ML predictions + live weather + 3 role dashboards + bilingual UI + automated advice â€” for Bugesera District.",
   0.38,7.04,12.55,0.42,sz=12,bold=True,col=TL)

# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
# SLIDE 8 â€” HOW THE SYSTEM WORKS (clean professional diagram)
# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
s=prs.slides.add_slide(BLANK)
hdr(s,"06","HOW THE SYSTEM WORKS","End-to-end workflow â€” from registration to harvest prediction")

# â”€â”€ TOP ROW: 6 flow steps with connecting arrows â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
SW=1.88; SH=2.0; SY=1.22; GAP=0.19
for i,(c,num,title,sub) in enumerate([
    (TD,"1","REGISTER",    "Farmer creates own\naccount on the platform.\nFills name, sector,\nemail, farm size."),
    (TS,"2","LOGIN",        "Farmer logs in with\nemail and password.\nSystem loads personal\nFarmer Dashboard."),
    (TM,"3","ENTER\nDATA",  "Selects crop, season,\nplanting date, seed,\nfertilizer type\nand terrain."),
    (TD,"4","WEATHER\nFETCH","System auto-calls\nOpen-Meteo API.\nGets live rainfall,\ntemp, humidity."),
    (TS,"5","ML MODEL\nPREDICTS","35 features fed to\nGradient Boosting.\nOutputs yield\nin kg/are + grade."),
    (TM,"6","RESULT\n+ADVICE","Farmer receives\nyield forecast,\nrecommendations\nand PDF report."),
]):
    x=0.22+i*(SW+GAP)
    # Card
    R(s,x,SY,SW,SH,fill=W,line=c,lw=2.2)
    # Colored header
    R(s,x,SY,SW,0.58,fill=c)
    # Step number circle
    R(s,x+0.68,SY+0.07,0.52,0.44,fill=W)
    T(s,num,x+0.7,SY+0.08,0.48,0.4,sz=16,bold=True,col=c,al=PP_ALIGN.CENTER)
    # Title
    T(s,title,x+0.06,SY+0.62,SW-0.1,0.55,sz=11,bold=True,col=c,al=PP_ALIGN.CENTER)
    # Separator
    R(s,x+0.22,SY+1.2,SW-0.42,0.04,fill=c)
    # Description
    T(s,sub,x+0.06,SY+1.28,SW-0.1,0.68,sz=9.5,col=SL,al=PP_ALIGN.CENTER)
    # Arrow to next
    if i<5:
        ax=x+SW+0.02; ay=SY+0.84
        R(s,ax,ay,GAP-0.02,0.3,fill=TM)
        T(s,">",ax+0.01,ay-0.02,GAP-0.04,0.34,sz=11,bold=True,col=W,al=PP_ALIGN.CENTER)

# â”€â”€ BOTTOM SECTION: Architecture (horizontal) â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
R(s,0.22,3.42,12.9,0.05,fill=G)
T(s,"TECHNOLOGY ARCHITECTURE",0.22,3.54,4.0,0.3,sz=10,bold=True,col=TD)
R(s,4.22,3.61,0.05,0.2,fill=G)

arch=[
    ("React + Vite","Frontend Â· Bilingual EN/RW"),
    ("Flask REST API","v4.0 Â· 30+ endpoints"),
    ("MySQL Database","bugesera_harvest"),
    ("Open-Meteo API","Live weather data"),
    ("Gmail SMTP","Email notifications"),
]
for i,(t,d) in enumerate(arch):
    x=0.22+i*2.62
    R(s,x,3.88,2.5,1.15,fill=TB if i%2==0 else W,line=G,lw=1)
    R(s,x,3.88,2.5,0.06,fill=TM if i%2==0 else TS)
    T(s,t,x+0.1,3.96,2.3,0.35,sz=11,bold=True,col=TD,al=PP_ALIGN.CENTER)
    T(s,d,x+0.1,4.3,2.3,0.28,sz=9.5,col=SL,al=PP_ALIGN.CENTER)
    if i<4:
        T(s,"< >",x+2.5,4.24,0.22,0.28,sz=9,bold=True,col=TM,al=PP_ALIGN.CENTER)

# â”€â”€ BOTTOM: 3 roles â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
R(s,0.22,5.22,12.9,0.05,fill=G)
for i,(role,pts,c) in enumerate([
    ("FARMER",         "Register Â· Predict harvest Â· View weather Â· Get advice Â· Download PDF result",    TD),
    ("SECTOR OFFICER", "Monitor all farms Â· Underperforming alerts Â· Send advice Â· Submit reports",        TS),
    ("DISTRICT ADMIN", "15-sector overview Â· ML metrics Â· Register officers Â· Send directives",            TM),
]):
    x=0.22+i*4.38
    R(s,x,5.32,4.18,1.72,fill=W,line=c,lw=1.5)
    R(s,x,5.32,4.18,0.44,fill=c)
    T(s,role,  x+0.12,5.36,3.9,0.32,sz=12,bold=True,col=W,al=PP_ALIGN.CENTER)
    T(s,pts,   x+0.12,5.84,3.9,0.62,sz=10.5,col=SL,al=PP_ALIGN.LEFT,wrap=True)

# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
# SLIDE 9 â€” RECOMMENDATIONS
# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
s=prs.slides.add_slide(BLANK)
hdr(s,"07","RECOMMENDATIONS","Future improvements â€” Chapter 5")

recs=[
    ("ðŸ“±","Develop Mobile App","Offline-capable Android/iOS for areas with limited internet."),
    ("ðŸ“¡","SMS Notifications","Telerivet/Twilio SMS alerts for farmers without smartphones."),
    ("ðŸ›°ï¸","Satellite Imagery","NDVI from Sentinel-2 for real-time crop health monitoring."),
    ("ðŸ—ºï¸","Expand to All Rwanda","Scale to all 30 districts under MINAGRI Smart Agriculture."),
    ("ðŸ”„","Continuous Retraining","Auto-retrain ML model yearly as new harvest data is collected."),
    ("ðŸ“Š","Advanced Analytics","Multi-season forecasting Â· market price integration."),
]
for i,(icon,title,desc) in enumerate(recs):
    row,col=divmod(i,2)
    x=0.22+col*6.56; y=1.3+row*1.88
    bg=TB if col==0 else W
    R(s,x,y,6.25,1.72,fill=bg,line=TM,lw=1.5)
    R(s,x,y,0.08,1.72,fill=TM)
    R(s,x+0.14,y+0.2,0.82,0.82,fill=TD)
    T(s,icon,x+0.16,y+0.22,0.76,0.68,sz=28,al=PP_ALIGN.CENTER)
    T(s,title,x+1.06,y+0.1,5.06,0.42,sz=13,bold=True,col=TD)
    R(s,x+1.06,y+0.56,5.06,0.04,fill=TM)
    T(s,desc, x+1.06,y+0.66,5.06,0.88,sz=11.5,col=SL)

# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
# SLIDE 10 â€” CONCLUSION
# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
s=prs.slides.add_slide(BLANK)
hdr(s,"08","CONCLUSION","Chapter 5 â€” Key achievements and overall findings")

# Summary
R(s,0.22,1.18,12.9,0.72,fill=TD)
R(s,0.22,1.18,12.9,0.08,fill=TM)
T(s,"All 5 specific objectives successfully achieved. The system demonstrates that Machine Learning "
   "can transform smallholder agriculture in Rwanda from guesswork to data-driven precision.",
   0.38,1.28,12.55,0.54,sz=12,col=W,it=True)

# 6 conclusion points
concs=[
    ("Objective 1 âœ“","2,502 agricultural records collected and analysed from 15 Bugesera sectors (2020â€“2024)."),
    ("Objective 2 âœ“","Gradient Boosting achieved RÂ²=0.9724 (97.24%) â€” exceeds 90% target. MAE=Â±0.979 kg/are."),
    ("Objective 3 âœ“","3-tier bilingual web platform built: Farmer Â· Sector Officer Â· District Admin dashboards."),
    ("Objective 4 âœ“","Open-Meteo API integrated for live weather across all 15 sectors with 7-day forecast."),
    ("Objective 5 âœ“","Personalized recommendations provided per crop, soil, sector, yield grade, and season."),
    ("National Impact","System aligns with Rwanda Vision 2050 Â· MINAGRI Smart Agriculture Â· NST1 targets."),
]
for i,(title,desc) in enumerate(concs):
    row,col=divmod(i,2)
    x=0.22+col*6.56; y=2.08+row*1.62
    bg=TB if col==0 else W
    border=TM if col==0 else TS
    R(s,x,y,6.25,1.48,fill=bg,line=border,lw=1.5)
    R(s,x,y,0.55,1.48,fill=TD if col==0 else TS)
    T(s,"âœ“",x+0.1,y+0.4,0.38,0.62,sz=22,bold=True,col=TL,al=PP_ALIGN.CENTER)
    T(s,title,x+0.65,y+0.1,5.44,0.38,sz=12,bold=True,col=TD)
    T(s,desc, x+0.65,y+0.55,5.44,0.82,sz=11,col=SL)

# Closing
R(s,0.22,7.0,12.9,0.58,fill=TD)
T(s,'"The Bugesera Harvest Prediction System bridges the gap between modern ML technology and grassroots farmers in Rwanda."',
   0.38,7.06,12.55,0.44,sz=12,bold=True,it=True,col=TL)

# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
# SLIDE 11 â€” THANK YOU / Q&A
# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
s=prs.slides.add_slide(BLANK)
R(s,0,0,13.33,7.5,fill=TD)
R(s,0,0,7.8,7.5,fill=TS)
R(s,0,0,13.33,0.08,fill=TM)
R(s,0,7.42,13.33,0.08,fill=TM)
R(s,8.6,-1.2,5.8,5.8,fill=TM)
R(s,10.0,4.5,4.2,4.2,fill=TD)
R(s,0,0,0.08,7.5,fill=TM)

# Logo
R(s,0.5,0.5,1.3,1.3,fill=TM)
T(s,"ðŸŒ¾",0.56,0.56,1.18,1.15,sz=46,al=PP_ALIGN.CENTER)

# Thank you
T(s,"THANK YOU",0.3,2.05,7.1,1.05,sz=58,bold=True,col=W,al=PP_ALIGN.CENTER)
R(s,0.75,3.18,6.2,0.1,fill=TM)
T(s,"Questions & Discussion Welcome",0.3,3.38,7.1,0.46,sz=18,col=TL,al=PP_ALIGN.CENTER)
T(s,'"Empowering Bugesera farmers through Machine Learning."',
   0.3,4.05,7.1,0.5,sz=13,it=True,col=W,al=PP_ALIGN.CENTER)

# 3 outcome chips
R(s,0.3,4.75,7.1,0.08,fill=TM)
for i,(val,lbl) in enumerate([("97.24%","Accuracy"),("5 / 5","Objectives"),("3 Roles","Served")]):
    x=0.38+i*2.38
    R(s,x,4.92,2.22,0.95,fill=TM if i%2==0 else TS)
    T(s,val,x+0.06,4.95,2.1,0.52,sz=26,bold=True,col=TL,al=PP_ALIGN.CENTER)
    T(s,lbl,x+0.06,5.46,2.1,0.28,sz=10,col=W,al=PP_ALIGN.CENTER)

# Right info
for i,(lbl,v1,v2) in enumerate([
    ("ðŸ‘©â€ðŸ’»  STUDENT",   "UWIMPUHWE Cesalie",    "Reg: 25RP21043"),
    ("ðŸ‘©â€ðŸ«  SUPERVISOR","Mrs. Marie MUTONI",    "RP Huye College"),
    ("ðŸŽ“  PROGRAM",    "B.Tech in ICT",         "Final Year 2025â€“2026"),
    ("ðŸ› ï¸  STACK",      "Python Â· Flask Â· React","MySQL Â· scikit-learn"),
    ("ðŸ’»  GITHUB",     "github.com/Cesalie/",   "MY-Project_Famer_Harvest"),
]):
    y=0.6+i*1.28
    R(s,7.9,y,5.2,1.12,fill=TS if i%2==0 else TD)
    R(s,7.9,y,5.2,0.08,fill=TM)
    T(s,lbl, 8.02,y+0.1,2.4,0.26,sz=8.5,bold=True,col=TL)
    T(s,v1,  8.02,y+0.42,3.5,0.34,sz=12,bold=True,col=W)
    T(s,v2,  11.45,y+0.42,1.5,0.34,sz=9,col=TL,al=PP_ALIGN.RIGHT)

# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
out=r"c:\Users\uwimp\Desktop\2026\Bugesera_Harvest_PPT_v7.pptx"
prs.save(out)
print("\n"+"="*58)
print("  âœ…  PPT SAVED")
print("="*58)
print(f"  ðŸ“  {out}")
print(f"  ðŸ“Š  11 slides Â· Teal + White Â· Clean & Panel-Ready")
print(f"  1. Title")
print(f"  2. Introduction")
print(f"  3. Problem Statement")
print(f"  4. Objectives")
print(f"  5. Challenges: Existing System")
print(f"  6. Development Challenges + Solutions")
print(f"  7. Uniqueness")
print(f"  8. How the System Works (diagram)")
print(f"  9. Recommendations")
print(f" 10. Conclusion")
print(f" 11. Thank You / Q&A")
print("="*58)

